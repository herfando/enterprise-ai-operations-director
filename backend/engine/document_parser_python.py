import os


def clean_dataframe(df):

    # Menghapus baris yang seluruh kolomnya kosong
    df = df.dropna(how="all")

    # Menghilangkan nilai NaN agar aman dikirim sebagai JSON
    df = df.fillna("")

    return df


def extract_excel_file(file_path):

    import pandas as pd

    df = pd.read_excel(file_path, engine="openpyxl")

    df = clean_dataframe(df)

    return {
        "type": "xlsx",
        "columns": list(df.columns),
        "data": df.to_dict(orient="records"),
    }


def extract_csv_file(file_path):

    import pandas as pd

    df = pd.read_csv(file_path, sep=None, engine="python")

    df = clean_dataframe(df)

    return {
        "type": "csv",
        "columns": list(df.columns),
        "data": df.to_dict(orient="records"),
    }


def extract_pptx_file(file_path):

    from pptx import Presentation

    prs = Presentation(file_path)

    slides = []

    for index, slide in enumerate(prs.slides, start=1):

        slide_content = {"slide": index, "texts": [], "tables": []}

        for shape in slide.shapes:

            if shape.has_text_frame:

                text = shape.text.strip()

                if text:
                    slide_content["texts"].append(text)

            if shape.has_table:

                table = []

                for row in shape.table.rows:

                    table.append([cell.text.strip() for cell in row.cells])

                slide_content["tables"].append(table)

        # notes powerpoint

        if slide.has_notes_slide:

            notes = slide.notes_slide.notes_text_frame.text.strip()

            if notes:

                slide_content["notes"] = notes

        slides.append(slide_content)

    return {"type": "pptx", "slides": slides}


def parse_python_document(file_path, filename):

    ext = os.path.splitext(filename)[1].lower()

    if ext == ".xlsx":

        return extract_excel_file(file_path)

    if ext == ".csv":

        return extract_csv_file(file_path)

    if ext == ".pptx":

        return extract_pptx_file(file_path)

    raise Exception(f"Python parser tidak support {ext}")

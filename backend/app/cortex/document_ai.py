import os
import json
import time
import re

from app.database.snowflake import get_snowflake_connection

STAGE_NAME = "@AI_OPERATIONS.AI_CONFIG.DOCUMENT_STAGE"


# =====================================================
# VALIDASI RULE-BASED UNTUK DATA TABULAR (XLSX / CSV)
# =====================================================
# Tujuan: gantiin AI_COMPLETE untuk kasus normal, karena data
# MO ini terstruktur (kolom jelas) jadi validasi & deteksi
# masalah bisa dilakukan instant pakai kode biasa, tanpa nunggu
# panggilan LLM yang lambat & mahal.
#
# AI_COMPLETE baru dipanggil (opsional) kalau ada baris yang
# ke-flag berisiko, dan itu pun cuma baris yang di-flag yang
# dikirim (bukan seluruh dokumen), supaya promptnya kecil.
#
# CATATAN: nama kolom di bawah ini disesuaikan dengan sample
# template "MARKETING ORDER (MO)" yang kamu kasih. Kalau nama
# kolom di file asli beda, sesuaikan KEYWORD_MAP di bawah.

KEYWORD_MAP = {
    "product_name": ["product name"],
    "buyer_name": ["buyer name"],
    "buyer_type": ["buyer type"],
    "order_qty": ["order quantity", "quantity"],
    "unit_price": ["unit selling price", "selling price"],
    "unit_cost": ["unit production cost", "production cost"],
    "mo_date": ["mo date"],
    "delivery_date": ["delivery date"],
}


# =====================================================
# 10 DEPARTEMEN & "TANDA TANGAN" DOKUMEN PER DEPARTEMEN
# =====================================================
# DEPARTMENTS = daftar resmi 10 departemen yang dipakai di seluruh
# sistem (dropdown upload, validasi department_match, dsb). Urutan &
# penulisan nama HARUS konsisten dengan yang dipakai di frontend/DB,
# karena string ini yang dibandingkan langsung ke input `department`.

DEPARTMENTS = [
    "Marketing",
    "PPIC",
    "Produksi",
    "QCQA",
    "Warehouse",
    "K3",
    "RnD",
    "Purchasing",
    "Maintenance",
    "Finance",
]


# DEPARTMENT_SIGNATURES = "ciri khas" dokumen tiap departemen, dipakai
# untuk DETEKSI OTOMATIS dokumen ini sebenernya punya departemen mana
# -- terlepas dari departemen apa yang dipilih user waktu upload.
#
# Ada 2 jenis ciri khas yang dicek:
#   - column_keywords : dipakai kalau dokumennya TABULAR (xlsx/csv).
#                        Formatnya sama seperti KEYWORD_MAP di atas:
#                        nama_field -> list keyword yang dicari di
#                        nama kolom (case-insensitive).
#   - text_keywords    : dipakai kalau dokumennya FREEFORM (hasil OCR
#                        dari pdf/gambar, atau isi pptx). Tinggal
#                        dicari kata kunci itu ada di teksnya atau
#                        nggak.
#
# STATUS SEKARANG: baru "Marketing" yang LENGKAP, karena baru itu
# yang ada contoh template MO-nya (KEYWORD_MAP di atas). 9 departemen
# lain masih PLACEHOLDER:
#   - text_keywords sudah diisi tebakan yang wajar (boleh disesuaikan
#     sendiri kalau kurang pas).
#   - column_keywords masih kosong ({}) sampai ada contoh template
#     dokumen resmi dari masing-masing departemen tersebut. Begitu
#     ada contohnya, isi dengan pola yang sama seperti "Marketing"
#     di bawah (lihat KEYWORD_MAP), lalu dia otomatis ikut terdeteksi
#     di detect_department_from_columns().

DEPARTMENT_SIGNATURES = {
    "Marketing": {
        "document_type": "Marketing Order (MO)",
        "column_keywords": KEYWORD_MAP,
        "text_keywords": [
            "marketing order",
            "buyer name",
            "unit selling price",
            "mo date",
        ],
    },
    "PPIC": {
        "document_type": None,  # TODO: isi nama dokumen resmi PPIC, mis. "Production Planning Schedule"
        "column_keywords": {},  # TODO: isi kalau sudah ada contoh template PPIC
        "text_keywords": [
            "ppic",
            "production planning",
            "perencanaan produksi",
            "inventory control",
        ],
    },
    "Produksi": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "produksi",
            "production report",
            "laporan produksi",
            "work order",
        ],
    },
    "QCQA": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "qc",
            "qa",
            "quality control",
            "quality assurance",
            "inspeksi kualitas",
        ],
    },
    "Warehouse": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "warehouse",
            "gudang",
            "stok barang",
            "inbound",
            "outbound",
        ],
    },
    "K3": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "k3",
            "kesehatan dan keselamatan kerja",
            "safety",
            "hse",
            "insiden kerja",
        ],
    },
    "RnD": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "rnd",
            "r&d",
            "research and development",
            "penelitian dan pengembangan",
        ],
    },
    "Purchasing": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "purchasing",
            "pembelian",
            "purchase order",
            "procurement",
        ],
    },
    "Maintenance": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "maintenance",
            "perawatan",
            "pemeliharaan",
            "preventive maintenance",
        ],
    },
    "Finance": {
        "document_type": None,  # TODO
        "column_keywords": {},  # TODO
        "text_keywords": [
            "finance",
            "keuangan",
            "invoice",
            "financial report",
            "laporan keuangan",
        ],
    },
}


def _normalize_department_name(name):
    """Samain format nama departemen biar perbandingan konsisten
    (case-insensitive, spasi rapih)."""

    return re.sub(r"\s+", " ", str(name or "")).strip().lower()


def detect_department_from_columns(df_columns):
    """Deteksi dokumen tabular ini kolom-kolomnya paling cocok sama
    departemen mana, berdasarkan DEPARTMENT_SIGNATURES.

    Return (nama_departemen, document_type, skor_cocok 0.0-1.0).
    Kalau nggak ada satupun signature yang cocok (termasuk kalau
    semua departemen selain Marketing masih kosong column_keywords
    -nya), return (None, None, 0.0)."""

    best_department = None
    best_doc_type = None
    best_score = 0.0

    for dept, sig in DEPARTMENT_SIGNATURES.items():
        col_map = sig.get("column_keywords") or {}
        if not col_map:
            continue

        matched = sum(1 for kws in col_map.values() if _find_column(df_columns, kws))
        score = matched / len(col_map)

        if score > best_score:
            best_score = score
            best_department = dept
            best_doc_type = sig.get("document_type")

    return best_department, best_doc_type, best_score


def detect_department_from_text(text_lower):
    """Deteksi dokumen freeform (hasil OCR pdf/gambar, atau isi
    pptx) ini kata-katanya paling banyak nyebut ciri khas departemen
    mana. Return (nama_departemen, jumlah_keyword_cocok)."""

    best_department = None
    best_score = 0

    for dept, sig in DEPARTMENT_SIGNATURES.items():
        keywords = sig.get("text_keywords") or []
        score = sum(1 for kw in keywords if kw in text_lower)

        if score > best_score:
            best_score = score
            best_department = dept

    return best_department, best_score


def build_department_mismatch_alert(uploaded_department, detected_department, evidence):
    """Bikin kalimat alert kalau departemen yang DIPILIH user waktu
    upload BEDA sama departemen yang TERDETEKSI dari isi dokumen.

    `evidence` = alasan singkat kenapa terdeteksi departemen itu,
    mis. "berisi data buyer, order quantity, dan unit selling price
    (ciri khas Marketing Order)"."""

    return (
        f"File yang Anda upload terdeteksi sebagai dokumen departemen "
        f"{detected_department}, karena isinya {evidence}. Anda upload "
        f"file ini ke departemen {uploaded_department}, padahal "
        f"seharusnya untuk departemen {detected_department}. Silakan "
        f"upload ulang file ini ke departemen {detected_department} "
        f"yang benar."
    )


def _find_column(df_columns, keywords):
    """Cari nama kolom asli berdasarkan keyword (case-insensitive,
    toleran terhadap newline/spasi ganda di header)."""

    for col in df_columns:
        normalized = re.sub(r"\s+", " ", str(col)).strip().lower()
        for kw in keywords:
            if kw in normalized:
                return col

    return None


def _to_number(value):
    """Parse angka yang mungkin ada koma ribuan (mis. '1,100')."""

    if value is None:
        return None

    try:
        return float(str(value).replace(",", "").strip())
    except (ValueError, TypeError):
        return None


def _to_date(value):
    import pandas as pd

    try:
        parsed = pd.to_datetime(value, dayfirst=True, errors="coerce")
        return None if pd.isna(parsed) else parsed
    except Exception:
        return None


def validate_mo_dataframe(df, department, timer=None):
    """Validasi tabel MO secara rule-based. Return dict dengan
    struktur mirip hasil AI_COMPLETE sebelumnya, plus daftar baris
    yang di-flag (kalau ada) untuk opsional dianalisis lebih lanjut
    oleh AI."""

    timer = timer or StepTimer()

    with timer.measure("rule_validation.check_rows"):

        columns = list(df.columns)
        col_map = {key: _find_column(columns, kws) for key, kws in KEYWORD_MAP.items()}

        missing_columns = [key for key, col in col_map.items() if col is None]

        missing_data = []
        flagged_rows = []

        if missing_columns:
            missing_data.append(
                f"Kolom tidak ditemukan di file: {', '.join(missing_columns)} "
                "(cek KEYWORD_MAP kalau nama kolom di file berbeda)"
            )

        qty_values = []
        if col_map["order_qty"]:
            for v in df[col_map["order_qty"]]:
                num = _to_number(v)
                if num is not None:
                    qty_values.append(num)

        qty_median = sorted(qty_values)[len(qty_values) // 2] if qty_values else None

        for idx, row in df.iterrows():
            row_no = idx + 1
            problems = []

            # --- data wajib kosong ---
            for key in [
                "product_name",
                "buyer_name",
                "order_qty",
                "unit_price",
                "unit_cost",
            ]:
                col = col_map.get(key)
                if col and (
                    row[col] is None or str(row[col]).strip() in ("", "nan", "NaN")
                ):
                    problems.append(f"{col} kosong")

            # --- harga jual < biaya produksi (rugi) ---
            if col_map["unit_price"] and col_map["unit_cost"]:
                price = _to_number(row[col_map["unit_price"]])
                cost = _to_number(row[col_map["unit_cost"]])
                if price is not None and cost is not None and price < cost:
                    problems.append(
                        f"Rugi: harga jual ({price:g}) < biaya produksi ({cost:g})"
                    )

            # --- kuantitas 0/negatif atau jauh di luar kewajaran ---
            if col_map["order_qty"]:
                qty = _to_number(row[col_map["order_qty"]])
                if qty is not None:
                    if qty <= 0:
                        problems.append("Kuantitas nol atau negatif")
                    elif qty_median and (qty > qty_median * 5 or qty < qty_median / 5):
                        problems.append(
                            f"Kuantitas ({qty:g}) jauh berbeda dari median baris lain ({qty_median:g})"
                        )

            # --- jadwal: delivery date vs MO date ---
            if col_map["mo_date"] and col_map["delivery_date"]:
                mo_date = _to_date(row[col_map["mo_date"]])
                delivery_date = _to_date(row[col_map["delivery_date"]])
                if mo_date is not None and delivery_date is not None:
                    selisih_hari = (delivery_date - mo_date).days
                    if selisih_hari < 0:
                        problems.append("Delivery date lebih awal dari MO date")
                    elif selisih_hari < 2:
                        problems.append(
                            f"Waktu produksi cuma {selisih_hari} hari, kemungkinan tidak realistis"
                        )

            if problems:
                flagged_rows.append({"row": row_no, "problems": problems})

        overall_risk = "none"
        if flagged_rows:
            overall_risk = "high" if len(flagged_rows) > len(df) * 0.2 else "medium"

        # --- DETEKSI DEPARTEMEN DARI ISI KOLOM (bukan cuma asumsi "Marketing") ---
        # Sekarang dicek ke semua 10 departemen lewat DEPARTMENT_SIGNATURES,
        # bukan di-hardcode "Marketing" seperti sebelumnya. Kalau di masa
        # depan tabel yg diupload lebih cocok signature departemen lain
        # (setelah column_keywords departemen lain diisi), ini otomatis
        # ke-detect sebagai departemen itu.
        detected_department, detected_doc_type, match_score = (
            detect_department_from_columns(columns)
        )

        # Kalau nggak ada satupun signature match dengan skor bagus (>0.5),
        # anggap tetap dokumen MO/Marketing selama kolom2 MO utama ketemu
        # (fallback ke perilaku lama supaya tidak salah flag dokumen MO asli
        # sebagai "tidak diketahui" hanya karena 1-2 kolom opsional kosong).
        if detected_department is None or match_score <= 0.5:
            if len(missing_columns) < len(KEYWORD_MAP):
                detected_department = "Marketing"
                detected_doc_type = "Marketing Order (MO)"
            else:
                detected_department = "Tidak diketahui"
                detected_doc_type = "Tidak diketahui"

        department_input_norm = _normalize_department_name(department)
        department_detected_norm = _normalize_department_name(detected_department)
        department_match = department_input_norm == department_detected_norm

        department_alert = None
        if not department_match and detected_department != "Tidak diketahui":
            evidence = (
                f"berisi data khas {detected_doc_type} seperti Buyer Name, "
                "Order Quantity, dan Unit Selling Price"
                if detected_department == "Marketing"
                else f"cocok dengan kolom/kata kunci khas departemen {detected_department}"
            )
            department_alert = build_department_mismatch_alert(
                uploaded_department=department,
                detected_department=detected_department,
                evidence=evidence,
            )

        return {
            "document_type": detected_doc_type,
            "department_detected": detected_department,
            "department_match": department_match,
            "department_alert": department_alert,
            "missing_data": missing_data,
            "risk_level": overall_risk,
            "flagged_rows": flagged_rows,
            "total_rows": len(df),
            "flagged_count": len(flagged_rows),
        }


def generate_recommendation_from_rules(rule_result):
    """PENGGANTI AI_COMPLETE untuk data tabular.

    Sebelumnya "problem" dan "recommended_action" dibikin oleh
    MISTRAL-LARGE2 lewat AI_COMPLETE (lambat, ~148 detik/panggilan).
    Sekarang dibikin dari template berbasis kategori masalah yang
    udah ke-detect di validate_mo_dataframe -- instant, konsisten,
    dan tujuannya sama: kasih tau apa masalahnya + apa yang harus
    dilakukan.
    """

    flagged_rows = rule_result["flagged_rows"]
    department_alert = rule_result.get("department_alert")

    if not flagged_rows:
        if department_alert:
            return {
                "problem": department_alert,
                "recommended_action": (
                    f"Upload ulang file ini ke departemen "
                    f"{rule_result['department_detected']} yang sesuai."
                ),
            }
        return {
            "problem": "Tidak ada masalah terdeteksi pada dokumen ini.",
            "recommended_action": "Lanjutkan proses normal, tidak ada tindakan tambahan diperlukan.",
        }

    # kelompokkan masalah per kategori biar rekomendasinya spesifik
    categories = {
        "rugi": [],
        "kosong": [],
        "kuantitas": [],
        "jadwal": [],
        "lainnya": [],
    }

    for item in flagged_rows:
        row_no = item["row"]
        for problem_text in item["problems"]:
            lower = problem_text.lower()
            if "rugi" in lower:
                categories["rugi"].append(row_no)
            elif "kosong" in lower:
                categories["kosong"].append(row_no)
            elif "kuantitas" in lower:
                categories["kuantitas"].append(row_no)
            elif "hari" in lower or "delivery" in lower:
                categories["jadwal"].append(row_no)
            else:
                categories["lainnya"].append(row_no)

    problem_parts = []
    action_parts = []

    if categories["rugi"]:
        rows = ", ".join(str(r) for r in sorted(set(categories["rugi"])))
        problem_parts.append(
            f"Baris {rows} berpotensi rugi (harga jual di bawah biaya produksi)"
        )
        action_parts.append(
            f"Review ulang harga jual untuk baris {rows} sebelum diproses"
        )

    if categories["kosong"]:
        rows = ", ".join(str(r) for r in sorted(set(categories["kosong"])))
        problem_parts.append(f"Baris {rows} punya data wajib yang kosong")
        action_parts.append(
            f"Lengkapi data yang kosong pada baris {rows} sebelum lanjut"
        )

    if categories["kuantitas"]:
        rows = ", ".join(str(r) for r in sorted(set(categories["kuantitas"])))
        problem_parts.append(f"Baris {rows} punya kuantitas yang tidak wajar")
        action_parts.append(
            f"Konfirmasi ulang jumlah pesanan pada baris {rows} ke buyer/sales"
        )

    if categories["jadwal"]:
        rows = ", ".join(str(r) for r in sorted(set(categories["jadwal"])))
        problem_parts.append(
            f"Baris {rows} punya jadwal produksi/pengiriman yang mepet atau tidak masuk akal"
        )
        action_parts.append(
            f"Koordinasi dengan tim produksi untuk cek kelayakan jadwal baris {rows}"
        )

    if categories["lainnya"]:
        rows = ", ".join(str(r) for r in sorted(set(categories["lainnya"])))
        problem_parts.append(f"Baris {rows} punya masalah lain yang perlu dicek manual")
        action_parts.append(f"Cek manual detail baris {rows}")

    problem_text = "; ".join(problem_parts) + "."
    action_text = "; ".join(action_parts) + "."

    if department_alert:
        problem_text = f"{department_alert} Selain itu: {problem_text}"
        action_text = (
            f"Upload ulang file ini ke departemen "
            f"{rule_result['department_detected']} yang sesuai. "
            f"Setelah itu: {action_text}"
        )

    return {
        "problem": problem_text,
        "recommended_action": action_text,
    }


def analyze_freeform_rules(content, department):
    """PENGGANTI AI_COMPLETE untuk dokumen non-tabular (hasil OCR
    dari PDF/gambar via AI_PARSE_DOCUMENT, atau isi PPTX).

    Karena teksnya bebas (bukan tabel terstruktur), validasinya
    lebih sederhana: cek departemen disebut di teksnya atau tidak,
    dan cek apakah kontennya kosong/gagal terbaca. Ini heuristik,
    bukan reasoning penuh seperti LLM -- tapi tujuannya sama dan
    instant, tanpa panggil AI_COMPLETE.
    """

    text = (content or "").strip()
    text_lower = text.lower()

    # --- deteksi departemen dari isi teks, dicek ke semua 10 departemen
    # lewat DEPARTMENT_SIGNATURES (bukan cuma list kecil hardcoded lagi) ---
    department_detected, keyword_hits = detect_department_from_text(text_lower)
    if department_detected is None:
        department_detected = "Tidak diketahui"

    department_input_norm = _normalize_department_name(department)
    department_detected_norm = _normalize_department_name(department_detected)

    # kalau nama departemen yg dipilih user sendiri kesebut literal di
    # teks, anggap match juga (jaga2 kalau text_keywords belum lengkap)
    department_match = (
        department_input_norm == department_detected_norm
        or department_input_norm in text_lower
    )

    missing_data = [] if text else ["Isi dokumen kosong atau gagal terbaca"]

    if not text:
        risk_level = "high"
        problem = (
            "Dokumen tidak terbaca atau kosong setelah diproses AI_PARSE_DOCUMENT."
        )
        recommended_action = "Cek ulang kualitas scan/foto, lalu upload ulang."
        department_alert = None
    elif not department_match and department_detected != "Tidak diketahui":
        risk_level = "medium"
        evidence = f"menyebut kata kunci khas departemen {department_detected}"
        department_alert = build_department_mismatch_alert(
            uploaded_department=department,
            detected_department=department_detected,
            evidence=evidence,
        )
        problem = department_alert
        recommended_action = (
            f"Upload ulang file ini ke departemen {department_detected} yang sesuai."
        )
    elif not department_match:
        risk_level = "medium"
        problem = f"Departemen '{department}' tidak ditemukan disebut di isi dokumen."
        recommended_action = (
            "Konfirmasi manual apakah dokumen ini memang untuk departemen tersebut."
        )
        department_alert = None
    else:
        risk_level = "none"
        problem = "Tidak ada masalah terdeteksi."
        recommended_action = "Lanjutkan proses normal."
        department_alert = None

    return {
        "document_type": "Dokumen (non-tabular)",
        "department_detected": department_detected,
        "department_match": department_match,
        "department_alert": department_alert,
        "missing_data": missing_data,
        "risk_level": risk_level,
        "problem": problem,
        "recommended_action": recommended_action,
    }


# =====================================================
# TIMING HELPER
# =====================================================
# Dipakai buat ngukur berapa lama tiap tahap jalan.
# Hasilnya dikumpulin di satu dict "timings" yang ikut
# dibalikin di response, jadi kelihatan langsung tanpa
# perlu buka log server.


class StepTimer:
    def __init__(self):
        self.timings = {}

    def measure(self, label):
        return _StepContext(self, label)


class _StepContext:
    def __init__(self, timer, label):
        self.timer = timer
        self.label = label

    def __enter__(self):
        self.start = time.perf_counter()
        print(f"[TIMING] mulai: {self.label}")
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        elapsed = time.perf_counter() - self.start
        self.timer.timings[self.label] = round(elapsed, 3)
        status = "GAGAL" if exc_type else "selesai"
        print(f"[TIMING] {status}: {self.label} -> {elapsed:.3f}s")
        return False  # jangan menelan exception


# =====================================================
# DETEKSI TIPE FILE ASLI (BUKAN CUMA DARI NAMA FILE)
# =====================================================


def detect_real_file_type(file_path):
    """
    Baca beberapa byte pertama file untuk mengenali tipe aslinya,
    supaya tidak salah parsing kalau ekstensi di nama file keliru
    (mis. file .csv yang sebenarnya .xlsx, atau sebaliknya).
    """

    with open(file_path, "rb") as f:
        header = f.read(8)

    # xlsx/pptx/docx semuanya adalah file ZIP (PK\x03\x04)
    if header.startswith(b"PK\x03\x04"):
        return "zip_office"  # xlsx atau pptx, dibedakan lagi dari nama file

    # xls lama (BIFF/OLE2)
    if header.startswith(b"\xd0\xcf\x11\xe0"):
        return "xls_legacy"

    # PDF
    if header.startswith(b"%PDF"):
        return "pdf"

    # PNG
    if header.startswith(b"\x89PNG"):
        return "png"

    # JPG
    if header.startswith(b"\xff\xd8\xff"):
        return "jpg"

    # kalau tidak match signature manapun, anggap teks/csv
    return "text_or_csv"


# =====================================================
# UPLOAD FILE KE SNOWFLAKE STAGE
# =====================================================


def upload_to_stage(file_path, timer=None):

    timer = timer or StepTimer()

    with timer.measure("upload_to_stage.connect"):
        conn = get_snowflake_connection()
        cursor = conn.cursor()

    try:
        absolute_path = os.path.abspath(file_path)

        print("UPLOAD FILE:", absolute_path)

        sql = """
        PUT %s
        @AI_OPERATIONS.AI_CONFIG.DOCUMENT_STAGE
        AUTO_COMPRESS=FALSE
        OVERWRITE=TRUE
        """

        with timer.measure("upload_to_stage.put_execute"):
            cursor.execute(sql, ("file://" + absolute_path,))
            result = cursor.fetchall()

        print("PUT RESULT:", result)

        return result

    finally:
        with timer.measure("upload_to_stage.close_connection"):
            cursor.close()
            conn.close()


# =====================================================
# CORTEX DOCUMENT AI
# PDF / IMAGE / SCAN TULISAN TANGAN
# =====================================================


def parse_document_with_cortex(filename, timer=None):

    timer = timer or StepTimer()

    with timer.measure("cortex_parse.connect"):
        conn = get_snowflake_connection()
        cursor = conn.cursor()

    try:
        safe_filename = filename.replace("'", "''")

        sql = f"""
        SELECT AI_PARSE_DOCUMENT(
            TO_FILE(
                '{STAGE_NAME}',
                '{safe_filename}'
            )
        )
        """

        with timer.measure("cortex_parse.AI_PARSE_DOCUMENT_execute"):
            cursor.execute(sql)
            result = cursor.fetchone()

        if result is None:
            return None

        parsed = result[0]

        print("PARSED TYPE :", type(parsed))
        print("PARSED VALUE :", parsed)

        return parsed

    finally:
        with timer.measure("cortex_parse.close_connection"):
            cursor.close()
            conn.close()


# =====================================================
# EXCEL XLSX PARSER
# =====================================================


def extract_excel_file(file_path, timer=None):

    import pandas as pd

    timer = timer or StepTimer()

    try:
        with timer.measure("extract_excel.read_excel"):
            df = pd.read_excel(file_path, engine="openpyxl")

        content = df.to_string(index=False)

        print("EXCEL CONTENT:")
        print(content)

        return content

    except ImportError as e:
        raise Exception(
            "Excel parsing failed: library 'openpyxl' belum "
            f"ter-install. Jalankan: pip install openpyxl. Detail: {e}"
        )

    except Exception as e:
        raise Exception(f"Excel parsing failed: {e}")


# =====================================================
# CSV PARSER (dengan fallback delimiter otomatis)
# =====================================================


def extract_csv_file(file_path, timer=None):

    import pandas as pd

    timer = timer or StepTimer()

    try:
        # sep=None + engine="python" bikin pandas otomatis
        # mendeteksi delimiter (koma, titik koma, tab, dll)
        # sehingga lebih tahan terhadap format CSV yang tidak konsisten.
        with timer.measure("extract_csv.read_csv"):
            df = pd.read_csv(file_path, sep=None, engine="python")

        content = df.to_string(index=False)

        print("CSV CONTENT:")
        print(content)

        return content

    except Exception as e:
        raise Exception(f"CSV parsing failed: {e}")


# =====================================================
# POWERPOINT PPTX PARSER
# =====================================================


def extract_pptx_file(file_path, timer=None):

    timer = timer or StepTimer()

    try:
        from pptx import Presentation
    except ImportError as e:
        raise Exception(
            "PowerPoint parsing failed: library 'python-pptx' belum "
            f"ter-install. Jalankan: pip install python-pptx. Detail: {e}"
        )

    try:
        with timer.measure("extract_pptx.parse_and_extract"):
            prs = Presentation(file_path)

            lines = []

            for slide_number, slide in enumerate(prs.slides, start=1):

                lines.append(f"--- SLIDE {slide_number} ---")

                for shape in slide.shapes:

                    # teks biasa (title, body, text box, dll)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in paragraph.runs).strip()

                            if text:
                                lines.append(text)

                    # teks di dalam tabel
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip(" |"):
                                lines.append(row_text)

                # speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        lines.append(f"[NOTES] {notes_text}")

            content = "\n".join(lines)

            print("PPTX CONTENT:")
            print(content)

            return content

    except Exception as e:
        raise Exception(f"PowerPoint parsing failed: {e}")


# =====================================================
# FILE ROUTER
# =====================================================


def extract_document_content(file_path, filename, timer=None):

    timer = timer or StepTimer()

    ext = os.path.splitext(filename)[1].lower().lstrip(".")
    real_type = detect_real_file_type(file_path)

    print("FILE NAME:", filename)
    print("EXTENSION (dari nama file):", ext)
    print("TIPE ASLI (dari isi file):", real_type)

    # --- validasi silang: ekstensi nama file vs isi file asli ---
    if ext == "xlsx" and real_type != "zip_office":
        raise Exception(
            f"File '{filename}' berekstensi .xlsx tapi isinya bukan file "
            f"Excel yang valid (terbaca sebagai '{real_type}'). "
            "Kemungkinan file corrupt, atau salah rename ekstensi."
        )

    if ext == "csv" and real_type == "zip_office":
        raise Exception(
            f"File '{filename}' berekstensi .csv tapi isinya sebenarnya "
            "file Office (xlsx/pptx/docx) yang di-zip. Cek lagi nama file "
            "atau proses upload-nya, sepertinya ekstensi tertukar."
        )

    if ext in ["pdf", "png", "jpg", "jpeg"]:

        parsed = parse_document_with_cortex(filename, timer=timer)

        if parsed is None:
            raise Exception("AI_PARSE_DOCUMENT returned NULL")

        if isinstance(parsed, str):
            try:
                parsed = json.loads(parsed)
            except json.JSONDecodeError:
                parsed = {"content": parsed}

        return parsed.get("content", "")

    elif ext in ["xlsx"]:

        return extract_excel_file(file_path, timer=timer)

    elif ext in ["csv"]:

        return extract_csv_file(file_path, timer=timer)

    elif ext in ["pptx"]:

        return extract_pptx_file(file_path, timer=timer)

    else:

        raise Exception(f"Unsupported file {ext}")


# =====================================================
# FILE EXCEL / CSV / PPT
# DIBACA DENGAN CORTEX
# =====================================================


def extract_structured_file_with_cortex(filename):

    conn = get_snowflake_connection()
    cursor = conn.cursor()

    try:
        prompt = f"""
YOU ARE AN AI DOCUMENT ANALYZER.

A manufacturing company uploaded this file:

{filename}

READ THIS DOCUMENT CONTENT.

EXTRACT ALL IMPORTANT DATA.

RETURN RAW TEXT JSON:

{{
"document_type":"",
"content":""
}}
"""

        sql = """
        SELECT AI_COMPLETE(
            'MISTRAL-LARGE2',
            %s
        )
        """

        cursor.execute(sql, (prompt,))

        result = cursor.fetchone()

        return result[0]

    finally:
        cursor.close()
        conn.close()


# =====================================================
# [TIDAK DIPAKAI LAGI] AI OPERATIONS DIRECTOR VIA AI_COMPLETE
# =====================================================
# Fungsi ini sudah TIDAK dipanggil di analyze_document().
# Diganti oleh validate_mo_dataframe() + generate_recommendation_from_rules()
# untuk xlsx/csv, dan analyze_freeform_rules() untuk pdf/gambar/pptx.
# Dibiarkan di sini cuma buat referensi/rollback kalau suatu saat
# ternyata butuh analisis AI lagi untuk kasus yang lebih kompleks.


def analyze_with_operations_director(
    document_text, department, timer=None, max_chars=8000
):

    timer = timer or StepTimer()

    # --- BATASI UKURAN DOKUMEN YANG DIKIRIM KE LLM ---
    # Prompt yang kepanjangan = waktu inference yang jauh lebih
    # lama (dan biaya lebih besar). Kalau dokumen jauh lebih
    # panjang dari max_chars, potong dan kasih catatan di prompt
    # bahwa data sudah dipotong, supaya model tidak "mengarang"
    # seolah itu keseluruhan data.
    original_length = len(document_text) if document_text else 0

    if original_length > max_chars:
        document_text = (
            document_text[:max_chars]
            + f"\n\n[... DATA DIPOTONG, total asli {original_length} karakter, "
            f"hanya {max_chars} karakter pertama yang dikirim ...]"
        )
        print(
            f"[TIMING] document_text dipotong: {original_length} -> "
            f"{max_chars} karakter"
        )

    print(
        "[TIMING] panjang document_text yang dikirim ke AI_COMPLETE:",
        len(document_text),
    )

    with timer.measure("operations_director.connect"):
        conn = get_snowflake_connection()
        cursor = conn.cursor()

    try:
        prompt = f"""
YOU ARE AN AI OPERATIONS DIRECTOR
IN A MANUFACTURING COMPANY.

DEPARTMENT:

{department}

ANALYZE THIS DOCUMENT.

RETURN JSON ONLY:

{{
"document_type":"",
"department_detected":"",
"department_match":"",

"missing_data":[],

"risk_level":"",
"problem":"",
"recommended_action":"",

"data":{{}}
}}

DOCUMENT:

{document_text}
"""

        sql = """
        SELECT AI_COMPLETE(
            'MISTRAL-LARGE2',
            %s
        )
        """

        with timer.measure("operations_director.AI_COMPLETE_execute"):
            cursor.execute(sql, (prompt,))
            result = cursor.fetchone()

        return result[0]

    finally:
        with timer.measure("operations_director.close_connection"):
            cursor.close()
            conn.close()


# =====================================================
# MAIN ORCHESTRATOR
# =====================================================


def analyze_document(file_path, filename, department):

    timer = StepTimer()
    overall_start = time.perf_counter()

    ext = os.path.splitext(filename)[1].lower().lstrip(".")

    # --- UPLOAD KE STAGE CUMA KALAU BENERAN BUTUH AI_PARSE_DOCUMENT ---
    # Stage upload cuma dipakai oleh Cortex OCR (pdf/png/jpg). xlsx/csv
    # dibaca langsung dari file_path pakai pandas, jadi TIDAK PERLU
    # upload ke Snowflake stage sama sekali -- itu murni overhead
    # yang kebuang percuma kalau tetap dijalankan.
    needs_stage_upload = ext in ["pdf", "png", "jpg", "jpeg"]

    if needs_stage_upload:
        upload_result = upload_to_stage(file_path, timer=timer)
        print(upload_result)
    else:
        print(f"[SKIP] upload_to_stage dilewati untuk .{ext} (tidak butuh Cortex OCR)")

    content = extract_document_content(file_path, filename, timer=timer)

    print("CONTENT TYPE:", type(content))
    print("CONTENT LENGTH (karakter):", len(content) if content else 0)
    print("CONTENT:", content)

    # --- SEMUA ANALISIS SEKARANG RULE-BASED, TIDAK ADA AI_COMPLETE ---
    if ext in ["xlsx", "csv"]:

        import pandas as pd

        with timer.measure("read_dataframe_for_validation"):
            if ext == "xlsx":
                df = pd.read_excel(file_path, engine="openpyxl")
            else:
                df = pd.read_csv(file_path, sep=None, engine="python")

        rule_result = validate_mo_dataframe(df, department, timer=timer)
        recommendation = generate_recommendation_from_rules(rule_result)

        decision = json.dumps({**rule_result, **recommendation}, ensure_ascii=False)

    else:
        # PDF/gambar/pptx: content sudah didapat dari AI_PARSE_DOCUMENT
        # (untuk OCR) atau python-pptx, lalu dianalisis rule-based juga
        with timer.measure("freeform_rule_validation"):
            freeform_result = analyze_freeform_rules(content, department)

        decision = json.dumps(freeform_result, ensure_ascii=False)

    overall_elapsed = round(time.perf_counter() - overall_start, 3)
    timer.timings["TOTAL"] = overall_elapsed

    # cetak ringkasan urut dari yang paling lama, biar langsung
    # kelihatan bottleneck-nya tanpa perlu scroll log
    print("\n===== RINGKASAN TIMING (detik) =====")
    for label, seconds in sorted(
        timer.timings.items(), key=lambda x: x[1], reverse=True
    ):
        pct = (seconds / overall_elapsed * 100) if overall_elapsed else 0
        print(f"{label:45s} {seconds:8.3f}s  ({pct:5.1f}%)")
    print("=====================================\n")

    return {
        "department": department,
        "document": content,
        "extracted": decision,
        "timings": timer.timings,
    }
